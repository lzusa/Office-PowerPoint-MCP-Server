"""
Content management utilities for PowerPoint MCP Server.
Functions for slides, text, images, tables, charts, and shapes.
"""
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from typing import Dict, List, Tuple, Optional, Any
import tempfile
import os
import base64


def add_slide(presentation: Presentation, layout_index: int = 1) -> Tuple:
    """
    Add a slide to the presentation.
    
    Args:
        presentation: The Presentation object
        layout_index: Index of the slide layout to use
        
    Returns:
        A tuple containing the slide and its layout
    """
    layout = presentation.slide_layouts[layout_index]
    slide = presentation.slides.add_slide(layout)
    return slide, layout


def get_slide_elements_minimal(slide, slide_index: int) -> Dict:
    """
    Lightweight element reader: returns ONLY element type, text, and position.
    Designed to minimize context/token consumption when inspecting slides.
    
    Args:
        slide: The slide object
        slide_index: Index of the slide
        
    Returns:
        Dictionary with slide dimensions and a list of elements.
        Each element contains: index, name, type, position, text (if any).
    """
    from pptx.util import Inches

    slide_width = round(float(slide.part.package.presentation_part.presentation.slide_width) / Inches(1).emu, 2)
    slide_height = round(float(slide.part.package.presentation_part.presentation.slide_height) / Inches(1).emu, 2)

    # Shape type classification
    SHAPE_TYPE_MAP = {
        'AUTO_SHAPE': 'shape',
        'CALLOUT': 'shape',
        'CHART': 'chart',
        'COMMENT': 'comment',
        'CANVAS': 'group',
        'DIAGRAM': 'diagram',
        'EMBEDDED_OLE_OBJECT': 'ole',
        'FORM': 'form',
        'FREEFORM': 'shape',
        'GROUP': 'group',
        'INK': 'ink',
        'INK_COMMENT': 'ink_comment',
        'LINE': 'line',
        'LINKED_OLE_OBJECT': 'ole',
        'LINKED_PICTURE': 'image',
        'MEDIA': 'media',
        'OLE_CONTROL_OBJECT': 'ole',
        'PICTURE': 'image',
        'PLACEHOLDER': 'placeholder',
        'SCRIPT_ANCHOR': 'script',
        'SHAPE': 'shape',
        'TABLE': 'table',
        'TEXT_BOX': 'textbox',
        'UNSPECIFIED': 'unknown',
    }

    elements = []
    for i, shape in enumerate(slide.shapes):
        left_in = round(shape.left / Inches(1).emu, 2) if shape.left else 0
        top_in = round(shape.top / Inches(1).emu, 2) if shape.top else 0
        width_in = round(shape.width / Inches(1).emu, 2) if shape.width else 0
        height_in = round(shape.height / Inches(1).emu, 2) if shape.height else 0

        stype = str(shape.shape_type)
        elem_type = SHAPE_TYPE_MAP.get(stype, 'shape')

        elem = {
            'index': i,
            'name': shape.name,
            'type': elem_type,
            'left': left_in,
            'top': top_in,
            'width': width_in,
            'height': height_in,
        }

        if hasattr(shape, 'text_frame') and shape.text_frame:
            text = shape.text_frame.text.strip()
            if text:
                elem['text'] = text if len(text) <= 200 else text[:200] + '...'
            # Placeholder extra info
            try:
                ptype = str(shape.placeholder_format.type)
                elem['ph_type'] = ptype
                elem['ph_idx'] = shape.placeholder_format.idx
            except Exception:
                pass
        elif hasattr(shape, 'image'):
            try:
                elem['image_info'] = f"{shape.image.content_type}"
            except Exception:
                pass
        elif hasattr(shape, 'table'):
            table = shape.table
            rows = len(table.rows)
            cols = len(table.columns)
            cells = []
            for r in range(rows):
                for c in range(cols):
                    t = table.cell(r, c).text.strip()
                    if t:
                        cells.append({'r': r, 'c': c, 'text': t})
            elem['table_size'] = f"{rows}x{cols}"
            if cells:
                elem['cells'] = cells

        elements.append(elem)

    return {
        'slide_index': slide_index,
        'slide_width': slide_width,
        'slide_height': slide_height,
        'element_count': len(elements),
        'elements': elements,
    }


def get_slide_info(slide, slide_index: int) -> Dict:
    """
    Get information about a specific slide.
    
    Args:
        slide: The slide object
        slide_index: Index of the slide
        
    Returns:
        Dictionary containing slide information
    """
    try:
        from pptx.util import Inches, Emu

        placeholders = []
        for placeholder in slide.placeholders:
            placeholder_info = {
                "idx": placeholder.placeholder_format.idx,
                "type": str(placeholder.placeholder_format.type),
                "name": placeholder.name
            }
            placeholders.append(placeholder_info)
        
        shapes = []
        for i, shape in enumerate(slide.shapes):
            left_emu = shape.left
            top_emu = shape.top
            width_emu = shape.width
            height_emu = shape.height

            left_in = float(Inches(1).emu / left_emu) if left_emu else 0
            top_in = float(Inches(1).emu / top_emu) if top_emu else 0
            width_in = float(width_emu) / Inches(1).emu if width_emu else 0
            height_in = float(height_emu) / Inches(1).emu if height_emu else 0

            shape_info = {
                "index": i,
                "name": shape.name,
                "shape_type": str(shape.shape_type),
                "left": left_emu,
                "top": top_emu,
                "width": width_emu,
                "height": height_emu,
                "left_inches": round(left_in, 2),
                "top_inches": round(top_in, 2),
                "width_inches": round(width_in, 2),
                "height_inches": round(height_in, 2),
                "has_text": bool(hasattr(shape, 'text_frame') and shape.text_frame and shape.text_frame.text.strip()),
            }
            shapes.append(shape_info)
        
        # Slide dimensions in inches (from presentation)
        try:
            pres = slide.part.package.presentation_part.presentation
            slide_width_emu = pres.slide_width
            slide_height_emu = pres.slide_height
        except Exception:
            slide_width_emu = None
            slide_height_emu = None
            
        slide_width_in = float(slide_width_emu) / Inches(1).emu if slide_width_emu else 0
        slide_height_in = float(slide_height_emu) / Inches(1).emu if slide_height_emu else 0

        return {
            "slide_index": slide_index,
            "layout_name": slide.slide_layout.name,
            "slide_width": slide_width_emu,
            "slide_height": slide_height_emu,
            "slide_width_inches": round(slide_width_in, 2),
            "slide_height_inches": round(slide_height_in, 2),
            "placeholder_count": len(placeholders),
            "placeholders": placeholders,
            "shape_count": len(shapes),
            "shapes": shapes
        }
    except Exception as e:
        raise Exception(f"Failed to get slide info: {str(e)}")


def set_title(slide, title: str) -> None:
    """
    Set the title of a slide.
    
    Args:
        slide: The slide object
        title: The title text
    """
    if slide.shapes.title:
        slide.shapes.title.text = title


def populate_placeholder(slide, placeholder_idx: int, text: str) -> None:
    """
    Populate a placeholder with text.
    
    Args:
        slide: The slide object
        placeholder_idx: The index of the placeholder
        text: The text to add
    """
    placeholder = slide.placeholders[placeholder_idx]
    placeholder.text = text


def add_bullet_points(placeholder, bullet_points: List[str]) -> None:
    """
    Add bullet points to a placeholder.
    
    Args:
        placeholder: The placeholder object
        bullet_points: List of bullet point texts
    """
    text_frame = placeholder.text_frame
    text_frame.clear()
    
    for i, point in enumerate(bullet_points):
        p = text_frame.add_paragraph()
        p.text = point
        p.level = 0


def add_textbox(slide, left: float, top: float, width: float, height: float, text: str,
                font_size: int = None, font_name: str = None, bold: bool = None,
                italic: bool = None, underline: bool = None, 
                color: Tuple[int, int, int] = None, bg_color: Tuple[int, int, int] = None,
                alignment: str = None, vertical_alignment: str = None, 
                auto_fit: bool = True) -> Any:
    """
    Add a textbox to a slide with formatting options.
    
    Args:
        slide: The slide object
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        text: Text content
        font_size: Font size in points
        font_name: Font name
        bold: Whether text should be bold
        italic: Whether text should be italic
        underline: Whether text should be underlined
        color: RGB color tuple (r, g, b)
        bg_color: Background RGB color tuple (r, g, b)
        alignment: Text alignment ('left', 'center', 'right', 'justify')
        vertical_alignment: Vertical alignment ('top', 'middle', 'bottom')
        auto_fit: Whether to auto-fit text
        
    Returns:
        The created textbox shape
    """
    textbox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    
    textbox.text_frame.text = text
    
    # Apply formatting if provided
    if any([font_size, font_name, bold, italic, underline, color, bg_color, alignment, vertical_alignment]):
        format_text_advanced(
            textbox.text_frame,
            font_size=font_size,
            font_name=font_name,
            bold=bold,
            italic=italic,
            underline=underline,
            color=color,
            bg_color=bg_color,
            alignment=alignment,
            vertical_alignment=vertical_alignment
        )
    
    return textbox


def format_text(text_frame, font_size: int = None, font_name: str = None, 
                bold: bool = None, italic: bool = None, color: Tuple[int, int, int] = None,
                alignment: str = None) -> None:
    """
    Format text in a text frame.
    
    Args:
        text_frame: The text frame to format
        font_size: Font size in points
        font_name: Font name
        bold: Whether text should be bold
        italic: Whether text should be italic
        color: RGB color tuple (r, g, b)
        alignment: Text alignment ('left', 'center', 'right', 'justify')
    """
    alignment_map = {
        'left': PP_ALIGN.LEFT,
        'center': PP_ALIGN.CENTER,
        'right': PP_ALIGN.RIGHT,
        'justify': PP_ALIGN.JUSTIFY
    }
    
    for paragraph in text_frame.paragraphs:
        if alignment and alignment in alignment_map:
            paragraph.alignment = alignment_map[alignment]
            
        for run in paragraph.runs:
            font = run.font
            
            if font_size is not None:
                font.size = Pt(font_size)
            if font_name is not None:
                font.name = font_name
            if bold is not None:
                font.bold = bold
            if italic is not None:
                font.italic = italic
            if color is not None:
                r, g, b = color
                font.color.rgb = RGBColor(r, g, b)


def format_text_advanced(text_frame, font_size: int = None, font_name: str = None, 
                        bold: bool = None, italic: bool = None, underline: bool = None,
                        color: Tuple[int, int, int] = None, bg_color: Tuple[int, int, int] = None,
                        alignment: str = None, vertical_alignment: str = None) -> Dict:
    """
    Advanced text formatting with comprehensive options.
    
    Args:
        text_frame: The text frame to format
        font_size: Font size in points
        font_name: Font name
        bold: Whether text should be bold
        italic: Whether text should be italic
        underline: Whether text should be underlined
        color: RGB color tuple (r, g, b)
        bg_color: Background RGB color tuple (r, g, b)
        alignment: Text alignment ('left', 'center', 'right', 'justify')
        vertical_alignment: Vertical alignment ('top', 'middle', 'bottom')
    
    Returns:
        Dictionary with formatting results
    """
    result = {
        'success': True,
        'warnings': []
    }
    
    try:
        alignment_map = {
            'left': PP_ALIGN.LEFT,
            'center': PP_ALIGN.CENTER,
            'right': PP_ALIGN.RIGHT,
            'justify': PP_ALIGN.JUSTIFY
        }

        vertical_alignment_map = {
            'top': MSO_VERTICAL_ANCHOR.TOP,
            'middle': MSO_VERTICAL_ANCHOR.MIDDLE,
            'bottom': MSO_VERTICAL_ANCHOR.BOTTOM
        }
        
        # Enable text wrapping
        text_frame.word_wrap = True

        if vertical_alignment and vertical_alignment in vertical_alignment_map:
            text_frame.vertical_anchor = vertical_alignment_map[vertical_alignment]
        
        # Apply formatting to all paragraphs and runs
        for paragraph in text_frame.paragraphs:
            if alignment and alignment in alignment_map:
                paragraph.alignment = alignment_map[alignment]
            
            for run in paragraph.runs:
                font = run.font
                
                if font_size is not None:
                    font.size = Pt(font_size)
                if font_name is not None:
                    font.name = font_name
                if bold is not None:
                    font.bold = bold
                if italic is not None:
                    font.italic = italic
                if underline is not None:
                    font.underline = underline
                if color is not None:
                    r, g, b = color
                    font.color.rgb = RGBColor(r, g, b)
                if bg_color is not None:
                    try:
                        r, g, b = bg_color
                        font.highlight_color.rgb = RGBColor(r, g, b)
                    except Exception:
                        pass
        
        return result
        
    except Exception as e:
        result['success'] = False
        result['error'] = str(e)
        return result


def add_image(slide, image_path: str, left: float, top: float, width: float = None, height: float = None) -> Any:
    """
    Add an image to a slide.
    
    Args:
        slide: The slide object
        image_path: Path to the image file
        left: Left position in inches
        top: Top position in inches
        width: Width in inches (optional)
        height: Height in inches (optional)
        
    Returns:
        The created image shape
    """
    if width is not None and height is not None:
        return slide.shapes.add_picture(
            image_path, Inches(left), Inches(top), Inches(width), Inches(height)
        )
    elif width is not None:
        return slide.shapes.add_picture(
            image_path, Inches(left), Inches(top), Inches(width)
        )
    elif height is not None:
        return slide.shapes.add_picture(
            image_path, Inches(left), Inches(top), height=Inches(height)
        )
    else:
        return slide.shapes.add_picture(
            image_path, Inches(left), Inches(top)
        )


def add_table(slide, rows: int, cols: int, left: float, top: float, width: float, height: float) -> Any:
    """
    Add a table to a slide.
    
    Args:
        slide: The slide object
        rows: Number of rows
        cols: Number of columns
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        
    Returns:
        The created table shape
    """
    return slide.shapes.add_table(
        rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)
    )


def format_table_cell(cell, font_size: int = None, font_name: str = None, 
                     bold: bool = None, italic: bool = None, 
                     color: Tuple[int, int, int] = None, bg_color: Tuple[int, int, int] = None,
                     alignment: str = None, vertical_alignment: str = None) -> None:
    """
    Format a table cell.
    
    Args:
        cell: The table cell object
        font_size: Font size in points
        font_name: Font name
        bold: Whether text should be bold
        italic: Whether text should be italic
        color: RGB color tuple (r, g, b)
        bg_color: Background RGB color tuple (r, g, b)
        alignment: Text alignment
        vertical_alignment: Vertical alignment
    """
    # Format text
    if any([font_size, font_name, bold, italic, color, alignment]):
        format_text_advanced(
            cell.text_frame,
            font_size=font_size,
            font_name=font_name,
            bold=bold,
            italic=italic,
            color=color,
            alignment=alignment
        )
    
    # Set background color
    if bg_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(*bg_color)


def add_chart(slide, chart_type: str, left: float, top: float, width: float, height: float,
              categories: List[str], series_names: List[str], series_values: List[List[float]]) -> Any:
    """
    Add a chart to a slide.
    
    Args:
        slide: The slide object
        chart_type: Type of chart ('column', 'bar', 'line', 'pie', etc.)
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        categories: List of category names
        series_names: List of series names
        series_values: List of value lists for each series
        
    Returns:
        The created chart object
    """
    # Map chart type names to enum values
    chart_type_map = {
        'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
        'stacked_column': XL_CHART_TYPE.COLUMN_STACKED,
        'bar': XL_CHART_TYPE.BAR_CLUSTERED,
        'stacked_bar': XL_CHART_TYPE.BAR_STACKED,
        'line': XL_CHART_TYPE.LINE,
        'line_markers': XL_CHART_TYPE.LINE_MARKERS,
        'pie': XL_CHART_TYPE.PIE,
        'doughnut': XL_CHART_TYPE.DOUGHNUT,
        'area': XL_CHART_TYPE.AREA,
        'stacked_area': XL_CHART_TYPE.AREA_STACKED,
        'scatter': XL_CHART_TYPE.XY_SCATTER,
        'radar': XL_CHART_TYPE.RADAR,
        'radar_markers': XL_CHART_TYPE.RADAR_MARKERS
    }
    
    xl_chart_type = chart_type_map.get(chart_type.lower(), XL_CHART_TYPE.COLUMN_CLUSTERED)
    
    # Create chart data
    chart_data = CategoryChartData()
    chart_data.categories = categories
    
    for i, series_name in enumerate(series_names):
        if i < len(series_values):
            chart_data.add_series(series_name, series_values[i])
    
    # Add chart to slide
    chart_shape = slide.shapes.add_chart(
        xl_chart_type, Inches(left), Inches(top), Inches(width), Inches(height), chart_data
    )
    
    return chart_shape.chart


def format_chart(chart, has_legend: bool = True, legend_position: str = 'right',
                has_data_labels: bool = False, title: str = None,
                x_axis_title: str = None, y_axis_title: str = None,
                color_scheme: str = None) -> None:
    """
    Format a chart with various options.
    
    Args:
        chart: The chart object
        has_legend: Whether to show legend
        legend_position: Position of legend ('right', 'top', 'bottom', 'left')
        has_data_labels: Whether to show data labels
        title: Chart title
        x_axis_title: X-axis title
        y_axis_title: Y-axis title
        color_scheme: Color scheme to apply
    """
    try:
        # Set chart title
        if title:
            chart.chart_title.text_frame.text = title
        
        # Configure legend
        if has_legend:
            chart.has_legend = True
            # Note: Legend position setting may vary by chart type
        else:
            chart.has_legend = False
        
        # Configure data labels
        if has_data_labels:
            for series in chart.series:
                series.has_data_labels = True
        
        # Set axis titles if available
        try:
            if x_axis_title and hasattr(chart, 'category_axis'):
                chart.category_axis.axis_title.text_frame.text = x_axis_title
            if y_axis_title and hasattr(chart, 'value_axis'):
                chart.value_axis.axis_title.text_frame.text = y_axis_title
        except:
            pass  # Axis titles may not be available for all chart types
            
    except Exception:
        pass  # Graceful degradation for chart formatting


def _alignment_to_string(alignment) -> Optional[str]:
    if alignment is None:
        return None

    alignment_map = {
        PP_ALIGN.LEFT: 'left',
        PP_ALIGN.CENTER: 'center',
        PP_ALIGN.RIGHT: 'right',
        PP_ALIGN.JUSTIFY: 'justify'
    }

    return alignment_map.get(alignment, str(alignment))


def _vertical_anchor_to_string(anchor) -> Optional[str]:
    if anchor is None:
        return None

    anchor_map = {
        MSO_VERTICAL_ANCHOR.TOP: 'top',
        MSO_VERTICAL_ANCHOR.MIDDLE: 'middle',
        MSO_VERTICAL_ANCHOR.BOTTOM: 'bottom'
    }

    return anchor_map.get(anchor, str(anchor))


def _rgb_to_list(rgb) -> Optional[List[int]]:
    if rgb is None:
        return None

    try:
        return [int(rgb[0]), int(rgb[1]), int(rgb[2])]
    except Exception:
        return None


def _safe_get_highlight_color(font) -> Optional[List[int]]:
    try:
        highlight = font.highlight_color
    except Exception:
        return None
    if highlight is None:
        return None

    try:
        return _rgb_to_list(highlight.rgb)
    except Exception:
        return None


def _is_placeholder(shape) -> bool:
    try:
        _ = shape.placeholder_format
        return True
    except Exception:
        return False


def _extract_text_frame_formatting(text_frame) -> Dict:
    formatting = {
        "word_wrap": getattr(text_frame, 'word_wrap', None),
        "vertical_alignment": _vertical_anchor_to_string(getattr(text_frame, 'vertical_anchor', None)),
        "paragraphs": []
    }

    try:
        paragraphs = list(text_frame.paragraphs)
    except Exception:
        return formatting

    for paragraph_index, paragraph in enumerate(paragraphs):
        paragraph_info = {
            "index": paragraph_index,
            "alignment": _alignment_to_string(getattr(paragraph, 'alignment', None)),
            "runs": []
        }

        try:
            runs = list(paragraph.runs)
        except Exception:
            runs = []

        for run_index, run in enumerate(runs):
            font = getattr(run, 'font', None)
            font_size = None
            if font is not None and getattr(font, 'size', None) is not None:
                try:
                    font_size = int(font.size.pt)
                except Exception:
                    font_size = None

            color = None
            if font is not None and getattr(font, 'color', None) is not None:
                try:
                    color = _rgb_to_list(font.color.rgb)
                except Exception:
                    color = None

            bg_color = _safe_get_highlight_color(font) if font is not None else None

            paragraph_info["runs"].append({
                "index": run_index,
                "text": getattr(run, 'text', ''),
                "font_name": font.name if font is not None else None,
                "font_size": font_size,
                "bold": font.bold if font is not None else None,
                "italic": font.italic if font is not None else None,
                "underline": font.underline if font is not None else None,
                "color": color,
                "bg_color": bg_color
            })

        formatting["paragraphs"].append(paragraph_info)

    return formatting


def extract_slide_text_content(slide) -> Dict:
    """
    Extract all text content from a slide including placeholders and text shapes.
    
    Args:
        slide: The slide object to extract text from
        
    Returns:
        Dictionary containing all text content organized by source type
    """
    try:
        text_content = {
            "slide_title": "",
            "placeholders": [],
            "text_shapes": [],
            "table_text": [],
            "table_cells": [],
            "all_text_combined": ""
        }
        
        all_texts = []
        
        # Extract title from slide if available
        if hasattr(slide, 'shapes') and hasattr(slide.shapes, 'title') and slide.shapes.title:
            try:
                title_text = slide.shapes.title.text_frame.text.strip()
                if title_text:
                    text_content["slide_title"] = title_text
                    all_texts.append(title_text)
            except:
                pass
        
        # Extract text from all shapes
        for i, shape in enumerate(slide.shapes):
            shape_text_info = {
                "shape_index": i,
                "shape_name": shape.name,
                "shape_type": str(shape.shape_type),
                "text": ""
            }
            
            try:
                # Check if shape has text frame
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        shape_text_info["formatting"] = _extract_text_frame_formatting(shape.text_frame)
                        shape_text_info["text"] = text
                        all_texts.append(text)
                        
                        # Categorize by shape type
                        if _is_placeholder(shape):
                            # This is a placeholder
                            placeholder_info = shape_text_info.copy()
                            placeholder_info["placeholder_type"] = str(shape.placeholder_format.type)
                            placeholder_info["placeholder_idx"] = shape.placeholder_format.idx
                            text_content["placeholders"].append(placeholder_info)
                        else:
                            # This is a regular text shape
                            text_content["text_shapes"].append(shape_text_info)
                
                # Extract text from tables
                elif hasattr(shape, 'table'):
                    table_texts = []
                    table_cell_details = []
                    table = shape.table
                    for row_idx, row in enumerate(table.rows):
                        row_texts = []
                        for col_idx, cell in enumerate(row.cells):
                            cell_text = cell.text_frame.text.strip()
                            if cell_text:
                                row_texts.append(cell_text)
                                all_texts.append(cell_text)
                                table_cell_details.append({
                                    "row": row_idx,
                                    "col": col_idx,
                                    "text": cell_text,
                                    "formatting": _extract_text_frame_formatting(cell.text_frame)
                                })
                        if row_texts:
                            table_texts.append({
                                "row": row_idx,
                                "cells": row_texts
                            })
                    
                    if table_texts:
                        text_content["table_text"].append({
                            "shape_index": i,
                            "shape_name": shape.name,
                            "table_content": table_texts
                        })

                    if table_cell_details:
                        text_content["table_cells"].append({
                            "shape_index": i,
                            "shape_name": shape.name,
                            "cells": table_cell_details
                        })
                        
            except Exception as e:
                # Skip shapes that can't be processed
                continue
        
        # Combine all text
        text_content["all_text_combined"] = "\n".join(all_texts)
        
        return {
            "success": True,
            "text_content": text_content,
            "total_text_shapes": len(text_content["placeholders"]) + len(text_content["text_shapes"]),
            "has_title": bool(text_content["slide_title"]),
            "has_tables": len(text_content["table_text"]) > 0
        }
        
    except Exception as e:
        return {
            "success": False,
            "error": f"Failed to extract text content: {str(e)}",
            "text_content": None
        }


# ============================================================
# Table Row/Column Manipulation
# ============================================================

def _safe_cell_formatting(cell) -> Dict:
    """
    Safely capture the current formatting of a table cell.

    Returns a dict with all preserved formatting values, or None for unset.
    """
    fmt = {}
    tf = cell.text_frame

    # Paragraph alignment
    p = tf.paragraphs[0] if tf.paragraphs else None
    fmt["alignment"] = p.alignment if p else None

    # Run-level font properties (from the first run)
    run = p.runs[0] if (p and p.runs) else None
    if run:
        font = run.font
        fmt["bold"] = font.bold
        fmt["italic"] = font.italic
        fmt["underline"] = font.underline
        fmt["font_size"] = font.size
        fmt["font_name"] = font.name
        fmt["font_color"] = font.color.rgb if font.color.type is not None else None
    else:
        fmt["bold"] = None
        fmt["italic"] = None
        fmt["underline"] = None
        fmt["font_size"] = None
        fmt["font_name"] = None
        fmt["font_color"] = None

    # Vertical alignment
    fmt["vertical_alignment"] = getattr(tf, "vertical_anchor", None)

    # Background fill
    fmt["bg_color"] = cell.fill.fore_color.rgb if cell.fill.type is not None else None

    return fmt


def _apply_cell_formatting(cell, fmt: Dict) -> None:
    """
    Apply a saved formatting dict back to a table cell.
    Only re-applies values that are not None.

    Works directly on the cell XML rather than going through format_table_cell,
    because format_table_cell expects int-point sizes and tuple-RGB colors,
    while font.size returns Centipoints and font.color.rgb returns RGBColor.
    """
    p = cell.text_frame.paragraphs[0] if cell.text_frame.paragraphs else None
    if p is None:
        return

    # Re-apply paragraph alignment
    if fmt.get("alignment") is not None:
        p.alignment = fmt["alignment"]

    # Re-apply run-level formatting
    for run in p.runs:
        font = run.font

        if fmt.get("font_size") is not None:
            font.size = fmt["font_size"]  # Centipoints object, assign directly
        if fmt.get("font_name") is not None:
            font.name = fmt["font_name"]
        if fmt.get("bold") is not None:
            font.bold = fmt["bold"]
        if fmt.get("italic") is not None:
            font.italic = fmt["italic"]
        if fmt.get("underline") is not None:
            font.underline = fmt["underline"]
        if fmt.get("font_color") is not None:
            font.color.rgb = fmt["font_color"]  # RGBColor object

    # Background fill
    if fmt.get("bg_color") is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fmt["bg_color"]

    # Vertical alignment
    if fmt.get("vertical_alignment") is not None:
        cell.text_frame.vertical_anchor = fmt["vertical_alignment"]


def edit_cell_text(cell, text: str = None, font_size: int = None, font_name: str = None,
                   bold: bool = None, italic: bool = None,
                   color: Tuple[int, int, int] = None, bg_color: Tuple[int, int, int] = None,
                   alignment: str = None, vertical_alignment: str = None) -> None:
    """
    Edit a table cell's text while preserving its existing formatting.

    Only the explicitly provided formatting parameters override the existing ones.

    Args:
        cell: The table cell object
        text: New text content for the cell (None to keep existing)
        font_size: Font size in points
        font_name: Font name
        bold: Whether text should be bold
        italic: Whether text should be italic
        color: RGB color tuple (r, g, b)
        bg_color: Background RGB color tuple (r, g, b)
        alignment: Text alignment
        vertical_alignment: Vertical alignment
    """
    if text is not None:
        # Preserve current formatting before replacing text
        saved_fmt = _safe_cell_formatting(cell)
        cell.text = str(text)
        _apply_cell_formatting(cell, saved_fmt)

    # Now apply any user-provided formatting overrides
    format_table_cell(
        cell,
        font_size=font_size,
        font_name=font_name,
        bold=bold,
        italic=italic,
        color=color,
        bg_color=bg_color,
        alignment=alignment,
        vertical_alignment=vertical_alignment
    )


def _create_tc(parent_tbl, default_height_emu: int = 400000) -> Any:
    """Create a new table cell XML element with an empty paragraph."""
    tc = parent_tbl.makeelement(qn('a:tc'), {})
    tc_pr = tc.makeelement(qn('a:tcPr'), {})
    tc.append(tc_pr)
    p = tc.makeelement(qn('a:p'), {})
    r = p.makeelement(qn('a:r'), {})
    rPr = r.makeelement(qn('a:rPr'), {})
    r.append(rPr)
    p.append(r)
    tc.append(p)
    return tc


def _create_grid_col(default_width: str = "914400") -> Any:
    """Create a new gridCol XML element."""
    return _create_grid_col_with_width(int(default_width))


def _create_grid_col_with_width(width_emu: int) -> Any:
    """Create a new gridCol XML element with specified width."""
    gc = _create_tc(None).makeelement(qn('a:gridCol'), {'w': str(width_emu)})
    return gc


def add_table_row(table, insert_at: int = None, data: Optional[List[str]] = None,
                  height: Optional[float] = None) -> int:
    """
    Add a new row to a table.

    Args:
        table: The table object
        insert_at: Row index to insert at (None = append at end)
        data: Optional list of text values for the new row cells
        height: Optional row height in inches (None = default)

    Returns:
        The index of the newly added row
    """
    tbl = table._tbl
    num_cols = len(table.columns)
    default_height_emu = int(height * 914400) if height else 400000

    # Create new <a:tr> element
    new_tr = tbl.makeelement(qn('a:tr'), {'h': str(default_height_emu)})

    # Add cells matching the grid columns
    for i in range(num_cols):
        tc = _create_tc(tbl)
        new_tr.append(tc)

    # Insert at position or append
    if insert_at is not None and insert_at < len(tbl.tr_lst):
        existing_tr = tbl.tr_lst[insert_at]
        tbl.insert(list(tbl.tr_lst).index(existing_tr), new_tr)
        row_index = insert_at
    else:
        tbl.append(new_tr)
        row_index = len(tbl.tr_lst) - 1

    table.notify_height_changed()

    # Populate data if provided
    if data:
        for i, value in enumerate(data):
            if i < num_cols:
                table.cell(row_index, i).text = str(value)

    return row_index


def delete_table_row(table, row_index: int) -> None:
    """
    Delete a row from a table by index.

    Args:
        table: The table object
        row_index: Index of the row to delete
    """
    tbl = table._tbl
    if row_index < 0 or row_index >= len(tbl.tr_lst):
        raise ValueError(f"Invalid row index: {row_index}. Available rows: 0-{len(tbl.tr_lst) - 1}")

    tr_to_remove = tbl.tr_lst[row_index]
    tbl.remove(tr_to_remove)
    table.notify_height_changed()


def add_table_column(table, insert_at: int = None, data: Optional[List[str]] = None,
                     width: Optional[float] = None) -> int:
    """
    Add a new column to a table.

    Args:
        table: The table object
        insert_at: Column index to insert at (None = append at end)
        data: Optional list of text values for the new column cells (one per row)
        width: Optional column width in inches (None = auto)

    Returns:
        The index of the newly added column
    """
    tbl = table._tbl
    num_rows = len(tbl.tr_lst)
    default_width_emu = int(width * 914400) if width else 914400

    # Create new <a:gridCol> element
    new_gc = tbl.tblGrid.makeelement(qn('a:gridCol'), {'w': str(default_width_emu)})

    # Insert at position or append
    if insert_at is not None and insert_at < len(tbl.tblGrid.gridCol_lst):
        existing_gc = list(tbl.tblGrid.gridCol_lst)[insert_at]
        tbl.tblGrid.insert(list(tbl.tblGrid.gridCol_lst).index(existing_gc), new_gc)
        col_index = insert_at
    else:
        tbl.tblGrid.append(new_gc)
        col_index = len(list(tbl.tblGrid.gridCol_lst)) - 1

    # Add a new <a:tc> to each existing row
    for tr in tbl.tr_lst:
        tc = _create_tc(tr)
        if insert_at is not None and insert_at < len(tr.findall(qn('a:tc'))):
            existing_tcs = tr.findall(qn('a:tc'))
            tr.insert(list(existing_tcs).index(existing_tcs[insert_at]), tc)
        else:
            tr.append(tc)

    table.notify_width_changed()

    # Populate data if provided
    if data:
        for i, value in enumerate(data):
            if i < num_rows:
                table.cell(i, col_index).text = str(value)

    return col_index


def delete_table_column(table, col_index: int) -> None:
    """
    Delete a column from a table by index.

    Args:
        table: The table object
        col_index: Index of the column to delete
    """
    tbl = table._tbl
    grid_cols = list(tbl.tblGrid.gridCol_lst)
    if col_index < 0 or col_index >= len(grid_cols):
        raise ValueError(f"Invalid column index: {col_index}. Available columns: 0-{len(grid_cols) - 1}")

    # Remove gridCol
    tbl.tblGrid.remove(grid_cols[col_index])

    # Remove corresponding tc from each row
    for tr in tbl.tr_lst:
        tcs = tr.findall(qn('a:tc'))
        if col_index < len(tcs):
            tr.remove(tcs[col_index])

    table.notify_width_changed()


def extract_slide_images(slide, slide_index: int, output_dir: str, naming_prefix: str = "slide") -> Dict:
    """
    Extract all images (pictures) from a slide and save them to disk.
    
    Image index matches get_slide_elements element index for easy cross-reference.

    Args:
        slide: The slide object
        slide_index: Index of the slide (used for naming)
        output_dir: Directory to save extracted images
        naming_prefix: Prefix for extracted image files (e.g. "slide" -> "slide_0_image_3.png")

    Returns:
        Dictionary with list of extracted image info or error.
    """
    import os

    os.makedirs(output_dir, exist_ok=True)

    images = []

    for i, shape in enumerate(slide.shapes):
        # Only picture shapes have embedded images
        if not hasattr(shape, 'image') or shape.image is None:
            continue

        try:
            img = shape.image
            content_type = img.content_type  # e.g. 'image/png', 'image/jpeg'

            # Determine file extension from content_type
            ext_map = {
                'image/png': '.png',
                'image/jpeg': '.jpg',
                'image/jpg': '.jpg',
                'image/gif': '.gif',
                'image/bmp': '.bmp',
                'image/tiff': '.tiff',
                'image/svg+xml': '.svg',
                'image/x-emf': '.emf',
                'image/x-wmf': '.wmf',
            }
            ext = ext_map.get(content_type, '.bin')

            # Use shape index (i) to match get_slide_elements element index
            file_name = f"{naming_prefix}_{slide_index}_image_{i}{ext}"
            file_path = os.path.join(output_dir, file_name)

            # Write image bytes
            with open(file_path, 'wb') as f:
                f.write(img.blob)

            images.append({
                'index': i,  # matches get_slide_elements element index
                'file_name': file_name,
                'file_path': file_path,
                'content_type': content_type,
                'size_bytes': len(img.blob),
                'shape_name': shape.name if hasattr(shape, 'name') else 'Unknown',
                'width_inches': round(shape.width / 914400, 2) if shape.width else None,
                'height_inches': round(shape.height / 914400, 2) if shape.height else None,
            })
        except Exception as e:
            images.append({
                'index': i,
                'error': str(e),
                'shape_name': shape.name if hasattr(shape, 'name') else 'Unknown',
            })

    return {
        'slide_index': slide_index,
        'image_count': len(images),
        'images': images,
    }


def extract_all_images(pres, output_dir: str, naming_prefix: str = "presentation") -> Dict:
    """
    Extract all images from every slide in a presentation.

    Args:
        pres: The Presentation object
        output_dir: Directory to save extracted images
        naming_prefix: Prefix for extracted image files

    Returns:
        Dictionary with summary and list of extracted images per slide.
    """
    results = []
    total_count = 0

    for idx, slide in enumerate(pres.slides):
        result = extract_slide_images(slide, idx, output_dir, naming_prefix)
        results.append(result)
        total_count += result.get('image_count', 0)

    return {
        'total_images': total_count,
        'slide_count': len(pres.slides),
        'output_dir': output_dir,
        'slides': results,
    }


# ── Label–Image matching & knowledge-base document builder ────────────

def _rect_distance(s1, s2) -> float:
    """
    Euclidean distance between two rectangular shapes (bounding boxes).
    Returns 0 when shapes overlap.
    All positions/sizes are in Emu (python-pptx native unit).
    """
    x1 = s1.left
    y1 = s1.top
    w1 = s1.width
    h1 = s1.height
    x2 = s2.left
    y2 = s2.top
    w2 = s2.width
    h2 = s2.height

    dx = max(0, max(x1, x2) - min(x1 + w1, x2 + w2))
    dy = max(0, max(y1, y2) - min(y1 + h1, y2 + h2))
    return (dx ** 2 + dy ** 2) ** 0.5


def _is_digit_label(shape) -> str | None:
    """Check if shape text is a pure digit (1, 2, ...). Returns digit string or None."""
    import os
    _debug = os.environ.get('DEBUG_GROUP_IMAGES')

    if not shape.has_text_frame:
        if _debug:
            print(f"  [_is_digit_label] SKIP (no text_frame) | shape_type={shape.shape_type}")
        return None
    text = shape.text_frame.text.strip()
    result = text if text.isdigit() else None
    if _debug:
        shape_desc = f"left={shape.left}, top={shape.top}, w={shape.width}, h={shape.height}"
        print(f"  [_is_digit_label] text={repr(text)} | isdigit={text.isdigit()} | result={repr(result)} | {shape_desc}")
    return result


def _is_potential_footnote(text_shape, img_shape) -> bool:
    """Check if text_shape is a small text box below img_shape (potential footnote/caption)."""
    if text_shape.width >= img_shape.width:
        return False
    if text_shape.height >= img_shape.height * 0.5:
        return False

    iy2 = img_shape.top + img_shape.height
    if text_shape.top < iy2:
        return False

    gap = text_shape.top - iy2
    if gap > img_shape.height * 2:
        return False

    tx1 = text_shape.left
    tx2 = text_shape.left + text_shape.width
    ix1 = img_shape.left
    ix2 = img_shape.left + img_shape.width

    if not (tx1 < ix2 and tx2 > ix1):
        return False

    return True


def _group_images_by_labels(slide, all_shapes=None):
    """
    Group images by their nearest label number.

    Two-phase strategy:
    1. **Phase One — labeled assignment**: each label claims its single
       nearest image (one-to-one, strict matching).
    2. **Phase Two — unlabeled inheritance**: every remaining unassigned
       image inherits the label of the closest already-grouped image
       (always inherits, no threshold guard — every image is returned).
       If the slide has *no* labels at all, all images are placed in
       a single unlabeled group (key `0`).

    Returns dict: {label_num: {'images': [...], 'footnotes': {img_id: footnote_text}}}
    """
    import os
    _debug = os.environ.get('DEBUG_GROUP_IMAGES')

    if all_shapes is None:
        all_shapes = list(slide.shapes)

    if _debug:
        print("\n" + "="*80)
        print("[DEBUG] _group_images_by_labels — Shape Classification")
        print("="*80)

    labels = []
    images = []
    text_shapes = []

    for shape in all_shapes:
        d = _is_digit_label(shape)
        if d is not None:
            labels.append((int(d), shape))
            if _debug:
                print(f"  [LABEL]  digit={d} | pos=({shape.left},{shape.top}) size=({shape.width}x{shape.height})")
        elif hasattr(shape, 'image') and shape.image is not None:
            images.append(shape)
            if _debug:
                print(f"  [IMAGE]  pos=({shape.left},{shape.top}) size=({shape.width}x{shape.height}) ct={shape.image.content_type if shape.image else 'N/A'}")
        elif shape.has_text_frame and shape.text_frame.text.strip():
            text_shapes.append(shape)
            if _debug:
                t = shape.text_frame.text.strip()[:60]
                print(f"  [TEXT]   text={repr(t)} | pos=({shape.left},{shape.top})")

    if _debug:
        print(f"\n  Summary: labels={len(labels)} | images={len(images)} | text_shapes={len(text_shapes)}")

    if not images:
        if _debug:
            print("  [WARN] No images found, returning empty dict.")
        return {}

    labels.sort(key=lambda x: x[0])
    has_labels = len(labels) > 0

    groups = {}           # label_num -> {'images': [...], 'footnotes': {}}
    image_labels = {}     # id(img) -> label_num  (tracks every image's label)

    # ── Phase 1: assign each label its single nearest image ──
    if has_labels:
        if _debug:
            print(f"\n{'='*80}")
            print(f"[DEBUG] Phase 1: Labeled Assignment (labels={[l[0] for l in labels]})")
            print(f"{'='*80}")

        used_images = set()

        # Compute threshold from all images
        threshold = max(img.width for img in images) * 1.5
        if _debug:
            print(f"  Distance threshold = max_img_width * 1.5 = {max(img.width for img in images)} * 1.5 = {threshold}")

        for label_num, lbl in labels:
            groups[label_num] = {'images': [], 'footnotes': {}}

            if _debug:
                print(f"\n  --- Label {label_num} (pos={lbl.left},{lbl.top}) ---")

            # Build candidate list from unassigned images
            candidates = [
                (img, _rect_distance(lbl, img))
                for img in images
                if id(img) not in used_images
            ]
            if _debug:
                if candidates:
                    for i, (img, dist) in enumerate(candidates):
                        status = f"{dist:.0f} -> {'WITHIN' if dist <= threshold else 'OVER'} threshold"
                        print(f"    Candidate {i}: img pos=({img.left},{img.top}) dist={dist:.0f} {status}")
                else:
                    print(f"    No candidates (all images already used)")

            if not candidates:
                if _debug:
                    print(f"    -> SKIP label {label_num}: no unassigned images left")
                continue                      # ← was `break` (fatal)

            candidates.sort(key=lambda x: x[1])

            # Claim the closest image
            closest_img, closest_dist = candidates[0]
            if _debug:
                print(f"    -> CLAIM closest: dist={closest_dist:.0f} img pos=({closest_img.left},{closest_img.top}) {'WITHIN threshold' if closest_dist <= threshold else '** OVER threshold **'}")
            groups[label_num]['images'].append(closest_img)
            used_images.add(id(closest_img))
            image_labels[id(closest_img)] = label_num

    # ── Phase 2: unlabeled image inheritance ──
    unassigned = [img for img in images if id(img) not in image_labels]

    if _debug:
        print(f"\n{'='*80}")
        print(f"[DEBUG] Phase 2: Unlabeled Image Inheritance")
        print(f"{'='*80}")
        print(f"  Unassigned images: {len(unassigned)}")
        for i, img in enumerate(unassigned):
            print(f"    Unassigned #{i}: pos=({img.left},{img.top}) size=({img.width}x{img.height})")

    if unassigned:
        if has_labels:
            # For each unassigned image, find the closest already-grouped
            # image and inherit its label.  Always inherit (no threshold
            # guard) so that every image is returned, even if the slide
            # layout is unusual.
            for img in unassigned:
                best_label = None
                best_dist = float('inf')
                best_gimg = None
                for label_num, group in groups.items():
                    for gimg in group['images']:
                        d = _rect_distance(img, gimg)
                        if d < best_dist:
                            best_dist = d
                            best_label = label_num
                            best_gimg = gimg

                if best_label is not None:
                    groups[best_label]['images'].append(img)
                    image_labels[id(img)] = best_label
                    if _debug:
                        print(f"  img pos=({img.left},{img.top}) -> INHERIT label {best_label} (dist={best_dist:.0f})")
                else:
                    # Should not happen if groups is non-empty, but fallback
                    first_label = next(iter(groups))
                    groups[first_label]['images'].append(img)
                    image_labels[id(img)] = first_label
                    if _debug:
                        print(f"  img pos=({img.left},{img.top}) -> FALLBACK to first_label {first_label}")
        else:
            # No labels at all: all images go into one unlabeled group
            groups[0] = {'images': list(images), 'footnotes': {}}
            for img in images:
                image_labels[id(img)] = 0
            if _debug:
                print(f"  No labels detected. All {len(images)} images grouped under key=0.")

    # ── Phase 3: footnote detection (unchanged) ──
    if _debug:
        print(f"\n{'='*80}")
        print(f"[DEBUG] Phase 3: Footnote Detection")
        print(f"{'='*80}")

    for label_num, group in groups.items():
        for img in group['images']:
            best_fn = None
            best_dist = float('inf')
            for ts in text_shapes:
                if _is_potential_footnote(ts, img):
                    d = _rect_distance(ts, img)
                    if d < best_dist:
                        best_dist = d
                        best_fn = ts
            if best_fn is not None:
                group['footnotes'][id(img)] = best_fn.text_frame.text.strip()
                if _debug:
                    fn_t = best_fn.text_frame.text.strip()[:50]
                    print(f"  Label {label_num} img pos=({img.left},{img.top}) -> footnote: {repr(fn_t)}")

    # ── Final summary ──
    if _debug:
        print(f"\n{'='*80}")
        print(f"[DEBUG] Final Groups Summary")
        print(f"{'='*80}")
        for label_num in sorted(groups.keys()):
            group = groups[label_num]
            img_count = len(group['images'])
            fn_count = len(group['footnotes'])
            print(f"  Group label={label_num}: {img_count} image(s), {fn_count} footnote(s)")
            for i, img in enumerate(group['images']):
                print(f"    img#{i} pos=({img.left},{img.top}) size=({img.width}x{img.height})")
        print(f"{'='*80}\n")

    return groups


def _extract_image_info(img_shape, footnote_text=None, label_num=None) -> dict:
    """Extract image metadata consistent with extract_images output format."""
    try:
        img = img_shape.image
        content_type = img.content_type
        ext_map = {
            'image/png': '.png', 'image/jpeg': '.jpg', 'image/jpg': '.jpg',
            'image/gif': '.gif', 'image/bmp': '.bmp', 'image/tiff': '.tiff',
            'image/svg+xml': '.svg', 'image/x-emf': '.emf', 'image/x-wmf': '.wmf',
        }
        ext = ext_map.get(content_type, '.bin')
        info = {
            'content_type': content_type,
            'size_bytes': len(img.blob),
            'width_inches': round(img_shape.width / 914400, 2) if img_shape.width else None,
            'height_inches': round(img_shape.height / 914400, 2) if img_shape.height else None,
        }
        if label_num is not None:
            info['label'] = label_num
        if footnote_text:
            info['footnote'] = footnote_text
        return info
    except Exception as e:
        return {'error': str(e)}


def build_slide_document(
    slide,
    slide_index: int,
    output_dir: str,
    naming_prefix: str = "slide",
) -> dict:
    """
    Build a knowledge-base document from a slide:
      1. Detect footnotes (small text boxes below images).
      2. Group images by label number (handles one label -> multiple images).
      3. Assign unlabeled images to nearest labeled image's group (rect-to-rect distance).
      4. If no labels at all, append all images at the end of text.
      5. Extract images to output_dir.
      6. Inject image references inline into text at (1), (2), ... markers.

    Returns:
        {
            'slide_index': int,
            'text': str,              # full text with inline image refs
            'images': [               # extracted image list
                {'label': 1, 'file_name': '...', 'file_path': '...', 'footnote': '...'},
            ],
            'image_count': int,
            'groups': [               # label grouping info
                {'label': 1, 'image_count': 2, 'has_footnote': True},
            ],
            'has_labels': bool,       # whether any labels were found
        }
    """
    import os, re

    _debug = os.environ.get('DEBUG_GROUP_IMAGES')

    os.makedirs(output_dir, exist_ok=True)

    ext_map = {
        'image/png': '.png', 'image/jpeg': '.jpg', 'image/jpg': '.jpg',
        'image/gif': '.gif', 'image/bmp': '.bmp', 'image/tiff': '.tiff',
        'image/svg+xml': '.svg', 'image/x-emf': '.emf', 'image/x-wmf': '.wmf',
    }

    # Step 1: group images by labels
    groups = _group_images_by_labels(slide)

    # Check if there are any real labels (not just label=0 for unlabeled)
    has_labels = any(k != 0 for k in groups.keys())

    # Collect footnote shape ids to skip them in text collection
    footnote_shape_ids = set()
    for group in groups.values():
        for img_id, fn_shape in group.get('footnotes', {}).items():
            if fn_shape is not None:
                footnote_shape_ids.add(id(fn_shape))

    # Step 2: extract images
    all_images = []  # list of (label_num, img_shape, file_path, info)
    label_to_images = {}  # label_num -> list of image info dicts

    image_counter = 0
    for label_num in sorted(groups.keys()):
        group = groups[label_num]
        label_to_images[label_num] = []
        for img_shape in group['images']:
            try:
                img = img_shape.image
                content_type = img.content_type
                ext = ext_map.get(content_type, '.bin')

                if has_labels and label_num != 0:
                    file_name = f"{naming_prefix}_label_{label_num}_{image_counter}{ext}"
                else:
                    file_name = f"{naming_prefix}_image_{image_counter}{ext}"

                file_path = os.path.join(output_dir, file_name)
                with open(file_path, 'wb') as f:
                    f.write(img.blob)

                footnote_text = None
                fn_shape = group['footnotes'].get(id(img_shape))
                if fn_shape is not None:
                    footnote_text = fn_shape.text_frame.text.strip()

                info = {
                    'label': label_num if has_labels and label_num != 0 else None,
                    'file_name': file_name,
                    'file_path': file_path,
                    'content_type': content_type,
                    'size_bytes': len(img.blob),
                    'width_inches': round(img_shape.width / 914400, 2) if img_shape.width else None,
                    'height_inches': round(img_shape.height / 914400, 2) if img_shape.height else None,
                }
                if footnote_text:
                    info['footnote'] = footnote_text

                all_images.append((label_num, img_shape, file_path, info))
                label_to_images[label_num].append(info)
                image_counter += 1
            except Exception as e:
                err_info = {'label': label_num if has_labels else None, 'error': str(e)}
                label_to_images[label_num].append(err_info)
                all_images.append((label_num, img_shape, None, err_info))

    # Step 3: collect text and inject image refs
    lines = []
    marker_pattern = re.compile(r'\((\d+)\)')
    has_markers = False

    # First pass: check if text contains any markers
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if text and marker_pattern.search(text):
            has_markers = True
            break

    if _debug:
        print(f"\n{'='*80}")
        print(f"[DEBUG] build_slide_document — Marker Replacement")
        print(f"{'='*80}")
        print(f"  has_markers={has_markers}")
        print(f"  label_to_images keys: {sorted(label_to_images.keys())}")
        for k, v in sorted(label_to_images.items()):
            print(f"    label_to_images[{k}]: {len(v)} image(s)")
            for vi in v:
                fn = vi.get('file_name', 'ERR')
                print(f"      -> {fn}")

    # Second pass: build text lines
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        # Skip footnote text boxes (already attached to images)
        if id(shape) in footnote_shape_ids:
            continue

        text = shape.text_frame.text.strip()
        if not text:
            continue

        paragraphs = text.split('\n')
        for para in paragraphs:
            para = para.strip()
            if not para:
                continue

            if has_markers:
                # Inject image refs at marker positions
                def _replace_marker(m, _label_to_images=label_to_images):
                    num = int(m.group(1))
                    imgs = _label_to_images.get(num)
                    if _debug:
                        if imgs:
                            print(f"  [MARKER] ({num}) -> FOUND {len(imgs)} image(s)")
                        else:
                            print(f"  [MARKER] ({num}) -> NOT FOUND (label_to_images has no key {num})")
                    if imgs:
                        refs = []
                        for img_info in imgs:
                            if 'error' in img_info:
                                continue
                            ref = f"[IMAGE: file={img_info['file_name']}, size={img_info['width_inches']}x{img_info['height_inches']}in, type={img_info['content_type']}]"
                            if img_info.get('footnote'):
                                ref += f" | footnote: {img_info['footnote']}"
                            refs.append(ref)
                        return f"{m.group()}{' '.join(refs)}"
                    return m.group()

                injected = marker_pattern.sub(_replace_marker, para)
                lines.append(injected)
            else:
                lines.append(para)

    full_text = '\n'.join(lines)

    # If no markers, append all images at the end
    if not has_markers and all_images:
        full_text += '\n\n--- 图片 ---\n'
        for label_num, img_shape, file_path, info in all_images:
            if 'error' in info:
                continue
            label_str = f"(标号{info['label']})" if info.get('label') is not None else "(无标号)"
            line = f"[IMAGE: file={info['file_name']}, size={info['width_inches']}x{info['height_inches']}in, type={info['content_type']}] {label_str}"
            if info.get('footnote'):
                line += f" | 脚注: {info['footnote']}"
            full_text += line + '\n'

    # Build groups summary
    groups_summary = []
    for label_num in sorted(groups.keys()):
        group = groups[label_num]
        has_footnote = any(fn is not None for fn in group['footnotes'].values())
        groups_summary.append({
            'label': label_num if has_labels and label_num != 0 else None,
            'image_count': len(group['images']),
            'has_footnote': has_footnote,
        })

    return {
        'slide_index': slide_index,
        'text': full_text,
        'images': [info for _, _, _, info in all_images if 'error' not in info],
        'image_count': len([info for _, _, _, info in all_images if 'error' not in info]),
        'groups': groups_summary,
        'has_labels': has_labels,
    }

